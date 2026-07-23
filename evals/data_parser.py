"""
Parses true_data (all files) and noisy_data (pptx/docx/txt only) into tagged chunks.
Uses python-docx and python-pptx directly — bypasses unstructured to avoid segfaults.
Reuses parse_text, parse_html, and chunk_text from the main app.
"""

import os      # Provides operating system utilities such as file and directory handling.
import sys     # Provides access to Python runtime and module search path.

# Add the project root directory to Python's module search path.
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import python-docx library for reading Word documents.
import docx as python_docx

# Import Presentation class for reading PowerPoint files.
from pptx import Presentation

# Import text parser for plain text and markdown files.
from app.ingestion.loaders.text import parse_text

# Import HTML parser.
from app.ingestion.loaders.html import parse_html

# Import chunking function that splits long documents into smaller chunks.
from app.ingestion.chunking.splitter import chunk_text

# Directory containing clean enterprise documents.
TRUE_DATA_DIR = os.path.join(
    os.path.dirname(os.path.dirname(__file__)),
    "data",
    "true_data"
)

# Directory containing noisy documents used during evaluation.
NOISY_DATA_DIR = os.path.join(
    os.path.dirname(os.path.dirname(__file__)),
    "data",
    "noisy_data"
)

# File types that are allowed inside noisy_data.
NOISY_ALLOWED_EXTS = {
    ".pptx",
    ".docx",
    ".txt"
}


# Reads a DOCX file and extracts all non-empty paragraphs.
def _parse_docx(file_path: str) -> str:

    # Open the Word document.
    doc = python_docx.Document(file_path)

    # Combine all non-empty paragraphs into one string.
    return "\n".join(
        p.text
        for p in doc.paragraphs
        if p.text.strip()
    )


# Reads a PowerPoint presentation and extracts text from every slide.
def _parse_pptx(file_path: str) -> str:

    # Open the PowerPoint presentation.
    prs = Presentation(file_path)

    # List that stores extracted text.
    texts = []

    # Iterate through every slide.
    for slide in prs.slides:

        # Iterate through every object on the slide.
        for shape in slide.shapes:

            # Check whether the object contains text.
            if hasattr(shape, "text") and shape.text.strip():

                # Store cleaned text.
                texts.append(shape.text.strip())

    # Merge all extracted text into a single string.
    return "\n".join(texts)


# Parses a file based on its extension.
def parse_file(file_path: str) -> str:

    # Extract file extension and convert it to lowercase.
    ext = os.path.splitext(file_path)[1].lower()

    try:

        # Parse Microsoft Word documents.
        if ext == ".docx":
            return _parse_docx(file_path)

        # Parse PowerPoint presentations.
        elif ext == ".pptx":
            return _parse_pptx(file_path)

        # Parse plain text or Markdown files.
        elif ext in (".txt", ".md"):
            return parse_text(file_path)

        # Parse HTML files.
        elif ext in (".html", ".htm"):
            return parse_html(file_path)

    # Ignore parsing errors and continue processing.
    except Exception:
        pass

    # Return an empty string if parsing failed or file type is unsupported.
    return ""


# Loads every document, parses it, chunks it, and tags each chunk.
def load_all_chunks() -> list[dict]:
    """
    Returns all chunks tagged with source filename and whether they are noise.
    Used by the eval pipeline to understand what context the RAG system draws from.
    """

    # Stores all generated chunks.
    results = []

    # ---------------- TRUE DATA ----------------

    # Iterate through every file in the true_data directory.
    for fname in sorted(os.listdir(TRUE_DATA_DIR)):

        # Construct the complete file path.
        fpath = os.path.join(TRUE_DATA_DIR, fname)

        # Skip directories.
        if not os.path.isfile(fpath):
            continue

        # Parse the document.
        text = parse_file(fpath)

        # Continue only if parsing succeeded.
        if text:

            # Split the document into chunks.
            for chunk in chunk_text(text):

                # Store chunk information.
                results.append({
                    "text": chunk,          # Chunk content.
                    "source": fname,        # Source filename.
                    "is_noise": False       # Marks it as clean data.
                })

    # ---------------- NOISY DATA ----------------

    # Iterate through every file in noisy_data.
    for fname in sorted(os.listdir(NOISY_DATA_DIR)):

        # Get file extension.
        ext = os.path.splitext(fname)[1].lower()

        # Ignore unsupported file types.
        if ext not in NOISY_ALLOWED_EXTS:
            continue

        # Build complete path.
        fpath = os.path.join(NOISY_DATA_DIR, fname)

        # Skip directories.
        if not os.path.isfile(fpath):
            continue

        # Parse the document.
        text = parse_file(fpath)

        # Continue only if parsing succeeded.
        if text:

            # Split into chunks.
            for chunk in chunk_text(text):

                # Store chunk information.
                results.append({
                    "text": chunk,         # Chunk content.
                    "source": fname,       # Source filename.
                    "is_noise": True       # Marks this chunk as noisy data.
                })

    # Return the complete list of chunks.
    return results